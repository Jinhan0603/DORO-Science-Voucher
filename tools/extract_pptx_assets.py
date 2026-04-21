from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _shape_text(shape: Any) -> str:
    if not hasattr(shape, "has_text_frame") or not shape.has_text_frame:
        return ""
    parts: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if text:
            parts.append(text)
    return " ".join(parts).strip()


def inspect_pptx(path: str | Path) -> dict[str, Any]:
    pptx_path = Path(path).resolve()
    prs = Presentation(str(pptx_path))

    slides: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        title = ""
        picture_count = 0
        media_count = 0
        picture_assets: list[dict[str, Any]] = []

        for shape_index, shape in enumerate(slide.shapes, start=1):
            text = _shape_text(shape)
            if text:
                texts.append(text)
                if not title:
                    title = text

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_count += 1
                try:
                    image = shape.image
                    picture_assets.append(
                        {
                            "shapeIndex": shape_index,
                            "name": image.filename,
                            "ext": image.ext,
                            "contentType": image.content_type,
                            "size": len(image.blob),
                        }
                    )
                except ValueError:
                    picture_assets.append(
                        {
                            "shapeIndex": shape_index,
                            "name": "",
                            "ext": "",
                            "contentType": "",
                            "size": 0,
                            "warning": "Picture shape has no embedded image payload.",
                        }
                    )
            elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                media_count += 1

        slide_text = " ".join(texts).strip()
        slides.append(
            {
                "slideNumber": slide_index,
                "title": title[:120] if title else "",
                "textSnippet": slide_text[:240],
                "pictureCount": picture_count,
                "mediaCount": media_count,
                "pictures": picture_assets,
            }
        )

    return {
        "path": str(pptx_path),
        "slideCount": len(prs.slides),
        "slides": slides,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description="Inspect a PPTX and return per-slide metadata.")
    parser.add_argument("pptx", help="Path to the PPTX file")
    parser.add_argument(
        "--json-out",
        help="Optional path to write JSON output. Prints to stdout when omitted.",
    )
    args = parser.parse_args()

    payload = inspect_pptx(args.pptx)
    output = json.dumps(payload, ensure_ascii=False, indent=2)

    if args.json_out:
        out_path = Path(args.json_out)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        out_path.write_text(output + "\n", encoding="utf-8")
    else:
        print(output)

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
