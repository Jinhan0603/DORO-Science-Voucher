from __future__ import annotations

import argparse
import io
import importlib.util
import json
import re
import shutil
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import fitz
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

LARGE_FILE_BYTES = 90 * 1024 * 1024
DEFAULT_IMAGE_MAX_SIDE = 1600
DEFAULT_JPEG_QUALITY = 82
RAW_IMAGE_MAX_SIDE = 1600

IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".svg", ".wdp"}
RASTER_IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff"}


@dataclass(frozen=True)
class KitSpec:
    kit_id: str
    kit_name_ko: str
    pdf_name: str
    pptx_name: str
    assembly_video_name: str
    operation_video_names: tuple[str, ...]
    extra_video_names: tuple[str, ...] = ()


KIT_SPECS: tuple[KitSpec, ...] = (
    KitSpec(
        kit_id="2-mood-light",
        kit_name_ko="스마트 무드등",
        pdf_name="1-1. 스마트 무드등_조립설명서.pdf",
        pptx_name="1-2. 스마트 무드등_조립설명서.pptx",
        assembly_video_name="1-3. 스마트 무드등_조립영상.mp4",
        operation_video_names=("무드등.mp4",),
        extra_video_names=("무드등 만드는과정.mp4",),
    ),
    KitSpec(
        kit_id="5-arduino-game",
        kit_name_ko="아두이노 게임기",
        pdf_name="2-1. 아두이노 게임기_조립설명서.pdf",
        pptx_name="2-2. 아두이노 게임기_조립설명서.pptx",
        assembly_video_name="2-3. 아두이노게임기_조립영상.mp4",
        operation_video_names=("게임기.mp4",),
    ),
    KitSpec(
        kit_id="1-bluetooth-speaker",
        kit_name_ko="블루투스 스피커",
        pdf_name="3-1. 블루투스 스피커_조립설명서.pdf",
        pptx_name="3-2. 블루투스 스피커_조립설명서.pptx",
        assembly_video_name="3-3. 블루투스 스피커_조립영상.mp4",
        operation_video_names=("스피커.mp4",),
    ),
    KitSpec(
        kit_id="3-walking-robot",
        kit_name_ko="워킹 로봇",
        pdf_name="4-1. 워킹로봇_조립설명서.pdf",
        pptx_name="4-2. 워킹로봇_조립설명서.pptx",
        assembly_video_name="4-3. 워킹로봇_조립영상 (2).mp4",
        operation_video_names=("워킹로봇.mp4",),
    ),
    KitSpec(
        kit_id="4-ir-car",
        kit_name_ko="IR 자동차",
        pdf_name="5-1. IR자동차_조립설명서.pdf",
        pptx_name="5-2. IR자동차_조립설명서.pptx",
        assembly_video_name="5-3. IR자동차_조립영상.mp4",
        operation_video_names=("IR자동차.mp4",),
    ),
    KitSpec(
        kit_id="6-ultrasonic-piano",
        kit_name_ko="초음파 피아노",
        pdf_name="6-1. 초음파피아노_조립설명서.pdf",
        pptx_name="6-2. 초음파피아노_조립설명서.pptx",
        assembly_video_name="6-3. 초음파피아노_조립영상.mp4",
        operation_video_names=("피아노.mp4",),
    ),
)

ASSEMBLY_SUBDIR = "6종_키트_안내서_조립PPT 및 영상"
DEMO_SUBDIR = "실제작동영상"
DEFAULT_OUTPUT_DIR = Path("assets/images/assembly")


def detect_tooling() -> dict[str, Any]:
    return {
        "python_pptx": importlib.util.find_spec("pptx") is not None,
        "pillow": importlib.util.find_spec("PIL") is not None,
        "pymupdf": importlib.util.find_spec("fitz") is not None,
        "ffmpeg": shutil.which("ffmpeg") is not None,
        "ffprobe": shutil.which("ffprobe") is not None,
        "magick": shutil.which("magick") is not None,
        "pdftoppm": shutil.which("pdftoppm") is not None,
        "soffice": shutil.which("soffice") is not None,
        "libreoffice": shutil.which("libreoffice") is not None,
        "powerpoint": shutil.which("powerpnt") is not None,
    }


def normalize_whitespace(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def safe_rel(path: Path, root: Path) -> str:
    return path.resolve().relative_to(root.resolve()).as_posix()


def file_record(path: Path | None, root: Path) -> dict[str, Any] | None:
    if not path or not path.exists():
        return None
    size = path.stat().st_size
    return {
        "path": safe_rel(path, root),
        "name": path.name,
        "size_bytes": size,
        "size_mb": round(size / 1024 / 1024, 2),
        "copy_to_repo_allowed": size < LARGE_FILE_BYTES,
    }


def find_named_file(folder: Path, filename: str) -> Path | None:
    candidate = folder / filename
    return candidate if candidate.exists() else None


def catalog_reference_materials(reference_root: Path) -> list[dict[str, Any]]:
    assembly_dir = reference_root / ASSEMBLY_SUBDIR
    demo_dir = reference_root / DEMO_SUBDIR
    kits: list[dict[str, Any]] = []

    for spec in KIT_SPECS:
        pdf_path = find_named_file(assembly_dir, spec.pdf_name)
        pptx_path = find_named_file(assembly_dir, spec.pptx_name)
        assembly_video = find_named_file(assembly_dir, spec.assembly_video_name)
        operation_videos = [p for name in spec.operation_video_names if (p := find_named_file(demo_dir, name))]
        extra_videos = [p for name in spec.extra_video_names if (p := find_named_file(demo_dir, name))]

        missing: list[str] = []
        if not pdf_path:
            missing.append("assembly_pdf")
        if not pptx_path:
            missing.append("assembly_pptx")
        if not assembly_video:
            missing.append("assembly_video")
        if not operation_videos:
            missing.append("operation_video")

        kits.append(
            {
                "kit_id": spec.kit_id,
                "kit_name_ko": spec.kit_name_ko,
                "assembly_pdf": pdf_path,
                "assembly_pptx": pptx_path,
                "assembly_video": assembly_video,
                "operation_videos": operation_videos,
                "extra_videos": extra_videos,
                "missing_assets": missing,
            }
        )

    return kits


def extract_text_from_shape(shape: Any) -> list[str]:
    texts: list[str] = []
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            texts.extend(extract_text_from_shape(child))
        return texts

    if getattr(shape, "has_text_frame", False):
        text = normalize_whitespace(shape.text)
        if text:
            texts.append(text)

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                text = normalize_whitespace(cell.text)
                if text:
                    texts.append(text)

    return texts


def collect_slide_texts(pptx_path: Path) -> list[dict[str, Any]]:
    prs = Presentation(str(pptx_path))
    slides: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        seen: set[str] = set()
        blocks: list[str] = []
        for shape in slide.shapes:
            for text in extract_text_from_shape(shape):
                if text not in seen:
                    seen.add(text)
                    blocks.append(text)
        slides.append(
            {
                "slide_number": idx,
                "text_blocks": blocks,
                "combined_text": " ".join(blocks),
            }
        )
    return slides


def infer_assembly_step_candidates(slides: list[dict[str, Any]]) -> list[int]:
    keywords = (
        "step",
        "단계",
        "조립",
        "연결",
        "끼우",
        "고정",
        "결합",
        "설치",
        "부착",
        "완성",
        "작동",
        "확인",
    )
    candidates: list[int] = []
    for slide in slides:
        combined = slide["combined_text"].lower()
        if any(keyword in combined for keyword in keywords):
            candidates.append(slide["slide_number"])
    return candidates


def ensure_clean_dir(path: Path) -> None:
    if path.exists():
        shutil.rmtree(path)
    path.mkdir(parents=True, exist_ok=True)


def render_pdf_slides(
    pdf_path: Path,
    output_dir: Path,
    *,
    max_side: int = DEFAULT_IMAGE_MAX_SIDE,
    jpeg_quality: int = DEFAULT_JPEG_QUALITY,
) -> list[Path]:
    ensure_clean_dir(output_dir)
    saved: list[Path] = []
    with fitz.open(str(pdf_path)) as doc:
        for idx, page in enumerate(doc, start=1):
            pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
            image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            longest_side = max(image.size)
            if longest_side > max_side:
                scale = max_side / float(longest_side)
                resized = (
                    max(1, int(image.width * scale)),
                    max(1, int(image.height * scale)),
                )
                image = image.resize(resized, Image.Resampling.LANCZOS)
            output_path = output_dir / f"step-{idx:02d}.jpg"
            image.save(output_path, "JPEG", quality=jpeg_quality, optimize=True, progressive=True)
            saved.append(output_path)
    return saved


def extract_pptx_media(pptx_path: Path, output_dir: Path) -> list[dict[str, Any]]:
    ensure_clean_dir(output_dir)
    extracted: list[dict[str, Any]] = []
    with zipfile.ZipFile(pptx_path) as archive:
        media_names = sorted(
            name
            for name in archive.namelist()
            if name.startswith("ppt/media/") and not name.endswith("/")
        )
        for idx, media_name in enumerate(media_names, start=1):
            suffix = Path(media_name).suffix.lower() or ".bin"
            if suffix not in IMAGE_SUFFIXES:
                continue

            output_path = output_dir / f"media-{idx:02d}{suffix}"
            data = archive.read(media_name)

            if suffix in RASTER_IMAGE_SUFFIXES:
                image = Image.open(io.BytesIO(data))
                if getattr(image, "is_animated", False):
                    image.seek(0)
                longest_side = max(image.size)
                if longest_side > RAW_IMAGE_MAX_SIDE:
                    scale = RAW_IMAGE_MAX_SIDE / float(longest_side)
                    resized = (
                        max(1, int(image.width * scale)),
                        max(1, int(image.height * scale)),
                    )
                    image = image.resize(resized, Image.Resampling.LANCZOS)

                if image.mode not in ("RGB", "RGBA"):
                    image = image.convert("RGBA" if "A" in image.getbands() else "RGB")

                if "A" in image.getbands() or suffix == ".png":
                    output_path = output_path.with_suffix(".png")
                    image.save(output_path, "PNG", optimize=True)
                else:
                    output_path = output_path.with_suffix(".jpg")
                    if image.mode != "RGB":
                        image = image.convert("RGB")
                    image.save(output_path, "JPEG", quality=DEFAULT_JPEG_QUALITY, optimize=True, progressive=True)
            else:
                with output_path.open("wb") as dst:
                    dst.write(data)

            extracted.append(
                {
                    "original_name": Path(media_name).name,
                    "saved_name": output_path.name,
                    "size_bytes": output_path.stat().st_size,
                }
            )
    return extracted


def extract_assets(
    reference_root: Path,
    repo_root: Path,
    output_root: Path | None = None,
) -> dict[str, Any]:
    output_root = output_root or (repo_root / DEFAULT_OUTPUT_DIR)
    kits = catalog_reference_materials(reference_root)
    tooling = detect_tooling()
    results: list[dict[str, Any]] = []

    for kit in kits:
        kit_root = output_root / kit["kit_id"]
        slides_dir = kit_root / "slides"
        raw_dir = kit_root / "raw"
        notes: list[str] = []

        slide_entries = collect_slide_texts(kit["assembly_pptx"]) if kit["assembly_pptx"] else []
        slide_images = render_pdf_slides(kit["assembly_pdf"], slides_dir) if kit["assembly_pdf"] else []
        raw_media = extract_pptx_media(kit["assembly_pptx"], raw_dir) if kit["assembly_pptx"] else []

        if kit["assembly_pptx"] and kit["assembly_pptx"].stat().st_size >= LARGE_FILE_BYTES:
            notes.append("PPTX 원본이 90MB 이상이라 repo에는 복사하지 않음")

        for video_path in [kit["assembly_video"], *kit["operation_videos"], *kit["extra_videos"]]:
            if video_path and video_path.stat().st_size >= LARGE_FILE_BYTES:
                notes.append(f"{video_path.name}은 90MB 이상으로 repo 직접 추가 금지")
            elif video_path:
                notes.append(f"{video_path.name}은 웹 반영 전 압축 권장")

        results.append(
            {
                "kit_id": kit["kit_id"],
                "kit_name_ko": kit["kit_name_ko"],
                "slide_count": len(slide_entries),
                "extractable_slide_image_count": len(slide_images),
                "raw_media_count": len(raw_media),
                "assembly_step_candidates": infer_assembly_step_candidates(slide_entries),
                "slides_dir": safe_rel(slides_dir, repo_root) if slides_dir.exists() else None,
                "raw_dir": safe_rel(raw_dir, repo_root) if raw_dir.exists() else None,
                "notes": notes,
            }
        )

    return {
        "reference_root": str(reference_root.resolve()),
        "repo_root": str(repo_root.resolve()),
        "output_root": safe_rel(output_root, repo_root),
        "tooling": tooling,
        "kits": results,
    }


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Extract slide images and embedded media from reference assets.")
    parser.add_argument(
        "--reference-root",
        type=Path,
        default=Path(__file__).resolve().parents[2] / "reference",
        help="Path to the external reference folder.",
    )
    parser.add_argument(
        "--repo-root",
        type=Path,
        default=Path(__file__).resolve().parents[1],
        help="Path to the site repository root.",
    )
    parser.add_argument(
        "--output-root",
        type=Path,
        default=None,
        help="Override output directory for extracted images.",
    )
    parser.add_argument(
        "--report-path",
        type=Path,
        default=None,
        help="Optional path to save an extraction summary JSON file.",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    result = extract_assets(
        reference_root=args.reference_root,
        repo_root=args.repo_root,
        output_root=args.output_root,
    )
    if args.report_path:
        args.report_path.parent.mkdir(parents=True, exist_ok=True)
        args.report_path.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
