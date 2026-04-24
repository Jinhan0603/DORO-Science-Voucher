from __future__ import annotations

import json
from pathlib import Path

import fitz
from PIL import Image
from pptx import Presentation

from assembly_source_config import ASSEMBLY_REFERENCE_DIR, KIT_CONFIGS, V7_ROOT


DATA_DIR = V7_ROOT / "assets" / "data"
PARTS_DIR = V7_ROOT / "assets" / "images" / "parts"
ASSEMBLY_DIR = V7_ROOT / "assets" / "images" / "assembly"
REPORT_PATH = DATA_DIR / "assembly-extraction-report.json"
TARGET_WIDTH = 1440
JPEG_QUALITY = 84


def ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def render_pdf_page(doc: fitz.Document, page_number: int) -> Image.Image:
    page = doc.load_page(page_number - 1)
    scale = TARGET_WIDTH / page.rect.width
    pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
    image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
    return image


def save_jpeg(image: Image.Image, path: Path) -> None:
    ensure_dir(path.parent)
    image.save(path, "JPEG", quality=JPEG_QUALITY, optimize=True, progressive=True)


def extract_kit_assets(kit_id: str, config: dict) -> dict:
    pdf_path = ASSEMBLY_REFERENCE_DIR / config["pdfName"]
    pptx_path = ASSEMBLY_REFERENCE_DIR / config["pptxName"]
    presentation = Presentation(str(pptx_path))
    pdf = fitz.open(str(pdf_path))

    parts_dir = PARTS_DIR / kit_id
    assembly_dir = ASSEMBLY_DIR / kit_id / "slides"
    ensure_dir(parts_dir)
    ensure_dir(assembly_dir)

    overview_path = parts_dir / "overview.jpg"
    overview_image = render_pdf_page(pdf, config["partsOverviewSlide"])
    save_jpeg(overview_image, overview_path)

    generated_steps = []
    for step in config["steps"]:
        step_file = assembly_dir / f"step-{int(step['step']):02d}.jpg"
        step_image = render_pdf_page(pdf, int(step["exportSlide"]))
        save_jpeg(step_image, step_file)
        generated_steps.append(
            {
                "step": step["step"],
                "sourceSlides": step["sourceSlides"],
                "exportSlide": step["exportSlide"],
                "asset": step_file.relative_to(V7_ROOT).as_posix(),
            }
        )

    slide_summaries = []
    for idx, slide in enumerate(presentation.slides, start=1):
        texts = []
        image_count = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:
                image_count += 1
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs).strip()
                    if line:
                        texts.append(line)
        slide_summaries.append(
            {
                "slide": idx,
                "images": image_count,
                "textPreview": " | ".join(texts[:6]),
            }
        )

    return {
        "pdf": pdf_path.name,
        "pptx": pptx_path.name,
        "slideCount": len(presentation.slides),
        "partsOverview": overview_path.relative_to(V7_ROOT).as_posix(),
        "generatedStepImages": generated_steps,
        "slideSummaries": slide_summaries,
    }


def main() -> None:
    ensure_dir(DATA_DIR)
    report = {}
    for kit_id, config in KIT_CONFIGS.items():
        report[kit_id] = extract_kit_assets(kit_id, config)

    REPORT_PATH.write_text(
        json.dumps(report, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    print(f"wrote {REPORT_PATH}")


if __name__ == "__main__":
    main()


